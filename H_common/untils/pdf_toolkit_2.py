import base64
import re
import os, io
from io import BytesIO
from pptx import Presentation
from hashlib import md5
import fitz
from PIL import Image
import imagehash
from PyPDF2 import PdfReader

from H_common.untils.files_utils import JsonUtils

"""
1、提取PPTX数据，转为JSON格式
2、提取pdf文本
3、提取pdf图片
4、提取pdf表格
"""


def pdf_table(file_path=''):
    """提取表格数据"""
    doc = fitz.open(file_path)
    results = {}
    # 确保保存 CSV 文件的目录存在
    save_dir = 'tables'
    os.makedirs(save_dir, exist_ok=True)
    for page in doc:
        tables = page.find_tables()
        num = page.number + 1
        num_list = []
        for i, table in enumerate(tables):
            df = table.to_pandas()
            csv_path = os.path.join(save_dir, f'table_pg_{num}_{i}.csv')
            df.to_csv(csv_path, index=False)
            num_list.append(csv_path)
        if num_list:
            results[num] = num_list
    return results

def pdf_text(file_path=''):
    """提取pdf文本"""
    reader = PdfReader(file_path)
    results = {}
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text = page.extract_text()
        results[page_num + 1] = text
    return results

def pdf_images(file_path=''):
    """
    提取pdf中的图片
    返回包含字典的列表，字典的key是pdf的页数，value是图片文件的存储全路径
    """
    pdf_document = fitz.open(file_path)
    seen_hashes = set()
    results = {}
    os.makedirs('images', exist_ok=True)
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        image_list = page.get_images(full=True)
        image_paths = []
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(io.BytesIO(image_bytes))
            image_hash = imagehash.average_hash(image)
            if image_hash not in seen_hashes:
                seen_hashes.add(image_hash)
                image_filename = f"images/page_{page_num + 1}_img_{img_index + 1}.png"
                image.save(image_filename)
                image_paths.append(image_filename)
                print(f"图片已保存为 {image_filename}")
        if image_paths:
            results[page_num + 1] = image_paths
    pdf_document.close()
    return results

def clean_text(text):
    cleaned_text = re.sub(r'\s+', ' ', text)
    cleaned_text = cleaned_text.strip()
    return cleaned_text

def image_to_base64(img):
    """图片转base64"""
    img_data = BytesIO()
    img.save(img_data, format=img.format)
    img_data.seek(0)
    base64_encoded_img = base64.b64encode(img_data.read()).decode('utf-8')
    return base64_encoded_img

def get_pdf_json_data(file_name, images, tables, texts):
    """获取pdf文件数据：JSON格式"""
    results = []
    for i, num in enumerate(texts):
        text = texts[num]
        table = tables[num] if num in tables else []
        image = images[num] if num in images else []
        res_num = {
            'file_name': file_name,
            "page": num,
            "text": text,
            "table_path": table,
            "image_path": image,
        }
        results.append(res_num)
    return results

def process_pptx(pptx_file):
    prs = Presentation(pptx_file)
    result_list = []
    result_dirt = {}

    for page_num, slide in enumerate(prs.slides, start=1):
        text_segments = ''
        image_data_list = []
        image_hashes = set()

        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text
                cleaned_text = clean_text(text)
                if cleaned_text:
                    text_segments += '\t' + cleaned_text

            if shape.shape_type == 13:
                image = shape.image
                image_stream = BytesIO(image.blob)
                img = Image.open(image_stream)

                img_hash = md5(image.blob).hexdigest()
                if img_hash not in image_hashes:
                    image_hashes.add(img_hash)
                    base64_encoded_img = image_to_base64(img)
                    image_data_list.append(base64_encoded_img)

        if text_segments.strip() or image_data_list:
            result = {
                "source_file": "袋鼠妈妈品牌介绍0618.pptx",
                "page": page_num,
                "text": text_segments.strip(),
                # "image": image_data_list
            }
            result_dirt[page_num] = result
            result_list.append(result)

    return result_list, result_dirt

def main():
    pptx_file = r'D:\D盘桌面\软通\袋鼠妈妈项目\知识库\AI营销数据\袋鼠妈妈品牌介绍0618.pptx'
    segment_file = r'D:\D盘桌面\软通\袋鼠妈妈项目\知识库\AI营销数据\袋鼠妈妈品牌介绍0618_原始拆分数据_20240813.json'
    jsonUtils = JsonUtils()
    result_list, _ = process_pptx(pptx_file)
    for page_result in result_list:
        print(page_result)
        jsonUtils.write_dict_to_json(page_result, segment_file)

def main2():
    pptx_file = r'D:\D盘桌面\软通\袋鼠妈妈项目\知识库\AI营销数据\袋鼠妈妈品牌介绍0618.pptx'
    _, result_dirt = process_pptx(pptx_file)

    txt_file = r'D:\D盘桌面\软通\袋鼠妈妈项目\知识库\AI营销数据\品宣数据-模型总结后的文本段.txt'
    line_count = 0

    segment_file = r'D:\D盘桌面\软通\袋鼠妈妈项目\知识库\AI营销数据\袋鼠妈妈品牌介绍0618_模型总结数据_20240813.json'

    jsonUtils = JsonUtils()
    with open(txt_file, 'r', encoding='utf-8') as file:
        for line in file:
            line_count += 1
            text = line.strip()
            print(text)
            if line_count == 1:
                page_num = [1, 2, 3, 4, 5]
            elif line_count == 2:
                page_num = [5]
            elif line_count == 3:
                page_num = [6, 7, 8]
            elif line_count == 4:
                page_num = [9, 10, 11, 12]
            elif line_count == 5:
                page_num = [8, 12]
            elif line_count == 6:
                page_num = [13, 14]
            elif line_count == 7:
                page_num = [15]
            elif line_count == 8:
                page_num = [16, 20]
            elif line_count == 9:
                page_num = [21, 22, 23, 24, 25]
            elif line_count == 10:
                page_num = [26, 27]

            text_segments_list = []
            for i, e in enumerate(page_num):
                data = result_dirt.get(e)
                if data:
                    text_segments_list.append(data['text'])
                if e == 7:
                    text_segments_list.append('在意每一个细节，宠爱每一个你\t安心\t关心\t专心\t用心\t累计3700万妈妈的选择。')

            result = {
                "source_file": "袋鼠妈妈品牌介绍0618.pptx",
                "text": text,
                "pptx_page": page_num,
                "pptx_text": text_segments_list
            }

            jsonUtils.write_dict_to_json(result, segment_file)

def main3():
    file_path = r"D:\D盘桌面\软通\袋鼠妈妈项目\知识库\测试数据\企业信用报告-广东袋鼠妈妈集团有限公司_20240730.pdf"
    file_name_with_ext = os.path.basename(file_path)
    file_name, _ = os.path.splitext(file_name_with_ext)
    result_json_file_path = f"D:\D盘桌面\软通\袋鼠妈妈项目\知识库\测试数据\{file_name}.json"
    # 提取表格数据
    tables = pdf_table(file_path)
    # 提取pdf图片
    images = pdf_images(file_path)
    # 提取pdf文本
    texts = pdf_text(file_path)
    # 获取pdf文件数据：JSON格式
    results = get_pdf_json_data(file_name_with_ext, images, tables, texts)
    # 写入json文件
    JsonUtils.write_list_to_json(results, result_json_file_path)

if __name__ == '__main__':
    # main()
    # main2()
    main3()

